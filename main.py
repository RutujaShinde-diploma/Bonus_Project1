from fastapi import FastAPI, File, Form, HTTPException, UploadFile
from fastapi.responses import FileResponse, HTMLResponse
from fastapi.middleware.cors import CORSMiddleware
import openai
import os
import tempfile
import shutil
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches
import json
import uuid

app = FastAPI(title="Text to PowerPoint Generator", version="1.0.0")

# Add CORS middleware
app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

# Temporary storage for generated files
TEMP_DIR = Path("temp_files")
TEMP_DIR.mkdir(exist_ok=True)

@app.post("/generate")
async def generate_presentation(
    input_text: str = Form(...),
    guidance: str = Form(""),
    api_key: str = Form(...),
    llm_provider: str = Form("openai"),
    template_file: UploadFile = File(...)
):
    """
    Generate a PowerPoint presentation from text input using a template
    """
    try:
        # Validate file type
        if not template_file.filename.lower().endswith(('.pptx', '.potx')):
            raise HTTPException(status_code=400, detail="Template must be .pptx or .potx file")
        
        # Create unique session ID
        session_id = str(uuid.uuid4())
        session_dir = TEMP_DIR / session_id
        session_dir.mkdir(exist_ok=True)
        
        # Save uploaded template
        template_path = session_dir / "template.pptx"
        with open(template_path, "wb") as buffer:
            shutil.copyfileobj(template_file.file, buffer)
        
        # Generate slide structure using LLM
        slide_structure = await generate_slide_structure(
            input_text, guidance, api_key, llm_provider
        )
        
        # Create PowerPoint presentation
        output_path = session_dir / "output.pptx"
        create_presentation(slide_structure, template_path, output_path)
        
        # Return file for download
        return FileResponse(
            path=output_path,
            filename="generated_presentation.pptx",
            media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation"
        )
        
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Error generating presentation: {str(e)}")
    finally:
        # Clean up temporary files after a delay (in production, use proper cleanup)
        pass

async def generate_slide_structure(text: str, guidance: str, api_key: str, provider: str) -> list:
    """
    Use LLM to generate slide structure from input text
    """
    if provider.lower() == "openai":
        return await generate_with_openai(text, guidance, api_key)
    else:
        # Default to OpenAI for now
        return await generate_with_openai(text, guidance, api_key)

async def generate_with_openai(text: str, guidance: str, api_key: str) -> list:
    """
    Generate slide structure using OpenAI API
    """
    try:
        client = openai.OpenAI(api_key=api_key)
        
        prompt = f"""
        Convert the following text into a PowerPoint presentation structure.
        
        Guidance: {guidance if guidance else 'Create a professional presentation'}
        
        Text: {text}
        
        Generate a JSON array of slides. Each slide should have:
        - title: slide title
        - content: array of bullet points or content
        
        Example format:
        [
            {{
                "title": "Introduction",
                "content": ["Welcome", "Agenda", "Key Objectives"]
            }},
            {{
                "title": "Main Content",
                "content": ["Point 1", "Point 2", "Point 3"]
            }}
        ]
        
        Create 3-8 slides based on the content. Make titles concise and content clear.
        Return only valid JSON, no other text.
        """
        
        response = client.chat.completions.create(
            model="gpt-3.5-turbo",
            messages=[{"role": "user", "content": prompt}],
            temperature=0.7,
            max_tokens=1000
        )
        
        content = response.choices[0].message.content.strip()
        
        # Try to extract JSON from response
        try:
            # Remove any markdown formatting
            if content.startswith("```json"):
                content = content[7:]
            if content.endswith("```"):
                content = content[:-3]
            
            slide_structure = json.loads(content.strip())
            return slide_structure
        except json.JSONDecodeError:
            # Fallback: create simple structure
            return create_fallback_structure(text)
            
    except Exception as e:
        print(f"OpenAI API error: {e}")
        return create_fallback_structure(text)

def create_fallback_structure(text: str) -> list:
    """
    Create a fallback slide structure if LLM fails
    """
    # Simple text splitting into slides
    words = text.split()
    words_per_slide = max(50, len(words) // 4)
    
    slides = []
    for i in range(0, len(words), words_per_slide):
        slide_words = words[i:i + words_per_slide]
        slide_text = " ".join(slide_words)
        
        # Create title from first few words
        title = " ".join(slide_words[:5]) + "..."
        if len(title) > 50:
            title = title[:47] + "..."
        
        slides.append({
            "title": title,
            "content": [slide_text]
        })
    
    return slides[:6]  # Limit to 6 slides

def create_presentation(slide_structure: list, template_path: Path, output_path: Path):
    """
    Create PowerPoint presentation using template and slide structure
    """
    try:
        # Load template
        prs = Presentation(template_path)
        
        # Get available layouts
        layouts = prs.slide_layouts
        
        # Find a suitable layout for content slides (usually index 1 or 2)
        content_layout = layouts[1] if len(layouts) > 1 else layouts[0]
        
        # Create new presentation or clear existing slides
        new_prs = Presentation(template_path)
        
        # Keep only the first slide (title slide) if it exists
        if len(new_prs.slides) > 0:
            # Remove all slides except the first
            while len(new_prs.slides) > 1:
                rId = new_prs.slides._sldIdLst[1].rId
                new_prs.part.drop_rel(rId)
                new_prs.slides._sldIdLst.remove(new_prs.slides._sldIdLst[1])
        else:
            # If no slides exist, create a title slide
            title_layout = layouts[0] if len(layouts) > 0 else content_layout
            slide = new_prs.slides.add_slide(title_layout)
            title = slide.shapes.title
            if title:
                title.text = "Generated Presentation"
        
        # Add content slides
        for slide_data in slide_structure:
            slide = new_prs.slides.add_slide(content_layout)
            
            # Set title
            if slide.shapes.title:
                slide.shapes.title.text = slide_data["title"]
            
            # Add content (try to find content placeholder)
            content_placeholder = None
            for shape in slide.placeholders:
                if shape.placeholder_format.type == 1:  # Content placeholder
                    content_placeholder = shape
                    break
            
            if content_placeholder and slide_data["content"]:
                # Add bullet points
                text_frame = content_placeholder.text_frame
                text_frame.clear()
                
                for i, point in enumerate(slide_data["content"]):
                    if i == 0:
                        p = text_frame.paragraphs[0]
                    else:
                        p = text_frame.add_paragraph()
                    
                    p.text = point
                    p.level = 0
        
        # Save presentation
        new_prs.save(output_path)
        
    except Exception as e:
        print(f"Error creating presentation: {e}")
        raise

@app.get("/")
async def root():
    return {"message": "Text to PowerPoint Generator API"}

@app.get("/index.html")
async def serve_frontend():
    """Serve the frontend HTML file"""
    try:
        with open("index.html", "r", encoding="utf-8") as f:
            html_content = f.read()
        return HTMLResponse(content=html_content)
    except FileNotFoundError:
        raise HTTPException(status_code=404, detail="Frontend not found")

@app.get("/health")
async def health_check():
    return {"status": "healthy"}

if __name__ == "__main__":
    import uvicorn
    uvicorn.run(app, host="0.0.0.0", port=8000)
